from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "AI-Powered Employee Skill Gap Analyser"
    subtitle1.text = "Intelligent Career Progression & Personalized Learning\n\nTeam D11 | Maitri · Ria · Sagar · Puneet"

    # --- Slide 2: Problem Statement ---
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    title2 = shapes2.title
    body2 = shapes2.placeholders[1]
    
    title2.text = "The Challenge in Modern Workforce Development"
    tf2 = body2.text_frame
    tf2.text = "Dynamic Job Roles: Technology evolves rapidly, making it hard to keep up."
    p = tf2.add_paragraph()
    p.text = "Vague Career Paths: Employees don't know exactly what skills they need to transition."
    p = tf2.add_paragraph()
    p.text = "Generic Training: HR prescribes 'one-size-fits-all' training instead of personalized plans."
    p = tf2.add_paragraph()
    p.text = "Manual Gap Analysis: Mapping skills against future requirements is tedious."

    # --- Slide 3: Our Solution ---
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    title3 = shapes3.title
    body3 = shapes3.placeholders[1]
    
    title3.text = "The Skill Gap Analyser App"
    tf3 = body3.text_frame
    tf3.text = "Deterministic Skill Assessment: Evaluates current proficiency against hardcoded industry standards."
    p = tf3.add_paragraph()
    p.text = "Instant Gap Matrix: Visually highlights areas of deficiency."
    p = tf3.add_paragraph()
    p.text = "AI Career Insights: Uses Groq LLM (Llama 3) for professional feedback."
    p = tf3.add_paragraph()
    p.text = "Targeted Recommendations: Queries Udemy dataset to suggest relevant courses."
    p = tf3.add_paragraph()
    p.text = "Auto-Generated Roadmap: Uses AI to generate a structured 6-month study plan."

    # --- Slide 4: System Architecture ---
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    title4 = shapes4.title
    body4 = shapes4.placeholders[1]
    
    title4.text = "How It Works Under The Hood"
    tf4 = body4.text_frame
    tf4.text = "Frontend: Streamlit dashboard for employee inputs."
    p = tf4.add_paragraph()
    p.text = "Core Engine: Standardizes inputs and calculates mathematical gap scores."
    p = tf4.add_paragraph()
    p.text = "Recommender: Parses missing skills against udemy_courses.csv."
    p = tf4.add_paragraph()
    p.text = "AI Integration: Connects to Groq Cloud (Llama 3) for roadmap generation."
    
    # Try to add the architecture image if it exists
    img_path = "architecture.png"
    if os.path.exists(img_path):
        slide4.shapes.add_picture(img_path, Inches(5), Inches(2.5), width=Inches(4.5))
        
    # --- Slide 5: Tech Stack ---
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    title5 = shapes5.title
    body5 = shapes5.placeholders[1]
    
    title5.text = "Technology Stack & Data Sources"
    tf5 = body5.text_frame
    tf5.text = "Language: Python"
    p = tf5.add_paragraph()
    p.text = "Frontend Framework: Streamlit (Custom CSS)"
    p = tf5.add_paragraph()
    p.text = "AI/LLM Provider: Groq API (Llama-3.3-70b)"
    p = tf5.add_paragraph()
    p.text = "Datasets: O*NET Database & Udemy Course Catalog"
    p = tf5.add_paragraph()
    p.text = "Reporting: Python-docx"

    # --- Slide 6: Future Scope ---
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    shapes6 = slide6.shapes
    title6 = shapes6.title
    body6 = shapes6.placeholders[1]
    
    title6.text = "Future Scope & Enhancements"
    tf6 = body6.text_frame
    tf6.text = "Live API Integrations: Pulling real-time course data from Udemy/Coursera."
    p = tf6.add_paragraph()
    p.text = "HR Dashboard: Unified view for managers to see department-wide gaps."
    p = tf6.add_paragraph()
    p.text = "Resume Parsing: Automatically extract current skill levels from CV uploads."
    p = tf6.add_paragraph()
    p.text = "Progress Tracking: Integrating interactive checkboxes for tracking milestones."

    # --- Slide 7: Conclusion ---
    slide7 = prs.slides.add_slide(title_slide_layout)
    title7 = slide7.shapes.title
    subtitle7 = slide7.placeholders[1]
    
    title7.text = "Empowering the Workforce of Tomorrow"
    subtitle7.text = "Bridging the gap between where employees are and where they want to be.\n\nThank You!\nAny Questions?"

    # Save presentation
    prs.save('Skill_Gap_Analyser_Presentation.pptx')
    print("Presentation saved as 'Skill_Gap_Analyser_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
